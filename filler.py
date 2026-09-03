"""Fill the BVIS proposal template with proposal data -> PPTX bytes."""
import io, os
from pptx import Presentation

_HERE = os.path.dirname(os.path.abspath(__file__))
_DEFAULT_TEMPLATE = os.path.join(_HERE, "template.pptx")

def _set(para, text):
    if not para.runs:
        r=para.add_run(); r.text=text; return
    para.runs[0].text=text
    for r in para.runs[1:]: r.text=""

def _fill_list(shape, header_text, items):
    """After a paragraph == header_text, fill following non-empty placeholder lines with items."""
    if not items: return
    paras=shape.text_frame.paragraphs
    started=False; idx=0
    for para in paras:
        t="".join(r.text for r in para.runs)
        if not started:
            if t.strip()==header_text: started=True
            continue
        if t.strip()=="" : continue
        # stop if we hit another bold header (heuristic: ends with ':')
        if t.strip().endswith(":") and idx>0: break
        if idx<len(items):
            _set(para, items[idx]); idx+=1
        else:
            _set(para, "")  # clear leftover template prompts

def fill(data: dict, template_path=None) -> bytes:
    p=Presentation(template_path or _DEFAULT_TEMPLATE)
    g=lambda k,d="": (data.get(k) or d)
    event=g("event","EVENT NAME"); committee=g("committee","Committee")
    month=g("month","Month + year"); year=g("year","2025 – 2026")
    footer=f"{data.get('event','Event Name')} – {committee}"
    acts=g("activities",[]) or []
    act1=g("act1_name", acts[0] if len(acts)>0 else "Activity 1")
    act2=g("act2_name", acts[1] if len(acts)>1 else "Activity 2")
    why=g("why",[]) or []
    allowed=g("allowed",[]) or []; notallowed=g("notallowed",[]) or []
    timeline=g("timeline",[]) or []

    SIMPLE={
      'PROPOSAL2025 – 2026 ACADEMIC YEAR': f'PROPOSAL{year} ACADEMIC YEAR',
      'Month + year':month,'Month + YEAR':month,'MONTH + YEAR':month.upper(),
      'Event Name – Name Committee':footer,
      'Date: ':f"Date: {g('date','TBC')}",'Location:':f"Location: {g('location','TBC')}",
      'Activity 1':act1,'Activity 2':act2,'Activity #2':act2,
      'ACTIVITY #1':act1.upper(),'ACTIVITY #2':act2.upper(),
    }
    for i,s in enumerate(p.slides):
        for sh in s.shapes:
            if not sh.has_text_frame: continue
            # simple paragraph replacements
            for para in sh.text_frame.paragraphs:
                t="".join(r.text for r in para.runs)
                if t in SIMPLE: _set(para, SIMPLE[t])
            # title red run
            if i==0:
                for para in sh.text_frame.paragraphs:
                    for r in para.runs:
                        if r.text.strip()=='EVENT NAME': r.text=event
            # overview lists (slide 2)
            if i==1:
                _fill_list(sh,'Why?',why)
            # activity1 detail (slide 4)
            if i==3:
                _fill_list(sh,'How students will participate:', g("act1_participate",[]) or [])
                _fill_list(sh,'Additional prizes and rewards:', g("act1_prizes",[]) or [])
            # activity2 detail (slide 8)
            if i==7:
                _fill_list(sh,'How students will participate:', g("act2_participate",[]) or [])
                _fill_list(sh,'Additional prizes and rewards:', g("act2_prizes",[]) or [])
            # timeline (slide 13)
            if i==12 and timeline:
                for para in sh.text_frame.paragraphs:
                    t="".join(r.text for r in para.runs)
                    if t.startswith('This is a timeline'):
                        _set(para, " · ".join(timeline))
    out=io.BytesIO(); p.save(out); return out.getvalue()

if __name__=="__main__":
    d=dict(event="HALLOWEEN HOUSE FEST",committee="Social Impact Committee",
      month="October 2026",year="2026 – 2027",date="30–31 October 2026",location="Auditorium",
      activities=["Pumpkin & Pins","Snatch the Donut"],
      why=["Suitable for all year groups","The donut game attracts everyone","Participants keep their reward; easy to set up"],
      act1_participate=["Students bowl a pumpkin at ghost-pins","Committee resets pins & records scores"],
      act1_prizes=["2 House Points per round winner","Postcard for top scorer"],
      timeline=["Wk1 buy materials","Wk2 promotion","Event day","Wk4 reflection"])
    open("out_test.pptx","wb").write(fill(d)); print("ok")
